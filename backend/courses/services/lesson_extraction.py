import io
from pathlib import Path

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


def _decode_text_bytes(raw_bytes):
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return raw_bytes.decode(encoding)
        except Exception:
            continue
    return raw_bytes.decode("utf-8", errors="ignore")


def _extract_from_pdf(raw_bytes):
    if fitz is None:
        return "", [], ["PyMuPDF is not installed. PDF text/image extraction skipped."]
    text_parts = []
    images = []
    warnings = []
    with fitz.open(stream=raw_bytes, filetype="pdf") as document:
        for page in document:
            text_parts.append(page.get_text("text") or "")
            for image_meta in page.get_images(full=True):
                xref = image_meta[0]
                extracted = document.extract_image(xref)
                if not extracted:
                    continue
                ext = extracted.get("ext", "png")
                blob = extracted.get("image")
                if not blob:
                    continue
                images.append({"name": f"pdf_page_{page.number + 1}_{xref}.{ext}", "bytes": blob})
    if not text_parts and not images:
        warnings.append("No extractable text/images found in PDF.")
    return "\n".join(text_parts).strip(), images, warnings


def _extract_from_docx(raw_bytes):
    if docx is None:
        return "", [], ["python-docx is not installed. DOCX extraction skipped."]
    text_parts = []
    images = []
    warnings = []
    document = docx.Document(io.BytesIO(raw_bytes))
    for paragraph in document.paragraphs:
        value = str(paragraph.text or "").strip()
        if value:
            text_parts.append(value)

    image_index = 0
    for relation in document.part.rels.values():
        if "image" not in relation.target_ref:
            continue
        image_index += 1
        blob = relation.target_part.blob
        suffix = Path(relation.target_ref).suffix or ".png"
        images.append({"name": f"docx_image_{image_index}{suffix}", "bytes": blob})

    if not text_parts and not images:
        warnings.append("No extractable text/images found in DOCX.")
    return "\n".join(text_parts).strip(), images, warnings


def _extract_from_pptx(raw_bytes):
    if Presentation is None:
        return "", [], ["python-pptx is not installed. PPTX extraction skipped."]
    text_parts = []
    images = []
    warnings = []
    deck = Presentation(io.BytesIO(raw_bytes))

    image_index = 0
    for slide in deck.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                value = str(shape.text_frame.text or "").strip()
                if value:
                    text_parts.append(value)
            if getattr(shape, "shape_type", None) == 13 and getattr(shape, "image", None):
                image_index += 1
                ext = shape.image.ext or "png"
                images.append({"name": f"pptx_image_{image_index}.{ext}", "bytes": shape.image.blob})

    if not text_parts and not images:
        warnings.append("No extractable text/images found in PPTX.")
    return "\n".join(text_parts).strip(), images, warnings


def extract_lesson_content(uploaded_file):
    """
    Returns:
      {
        "title_suggestion": str,
        "extracted_text": str,
        "images": [{"name": str, "bytes": bytes}],
        "warnings": [str],
        "extension": str,
      }
    """
    original_name = str(getattr(uploaded_file, "name", "") or "").strip()
    stem = Path(original_name).stem or "Imported Lesson"
    extension = Path(original_name).suffix.lower()
    raw_bytes = uploaded_file.read()
    if hasattr(uploaded_file, "seek"):
        uploaded_file.seek(0)

    if extension == ".pdf":
        extracted_text, images, warnings = _extract_from_pdf(raw_bytes)
    elif extension == ".docx":
        extracted_text, images, warnings = _extract_from_docx(raw_bytes)
    elif extension == ".pptx":
        extracted_text, images, warnings = _extract_from_pptx(raw_bytes)
    elif extension == ".txt":
        extracted_text = _decode_text_bytes(raw_bytes).strip()
        images = []
        warnings = []
    elif extension in {".doc", ".ppt"}:
        extracted_text = ""
        images = []
        warnings = [f"{extension} extraction is limited. File was accepted but no auto-extraction was performed."]
    else:
        extracted_text = ""
        images = []
        warnings = ["Unsupported file type for extraction."]

    return {
        "title_suggestion": stem.replace("_", " ").replace("-", " ").strip() or "Imported Lesson",
        "extracted_text": extracted_text,
        "images": images,
        "warnings": warnings,
        "extension": extension,
    }
