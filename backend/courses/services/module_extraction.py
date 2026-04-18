from __future__ import annotations

import io
import statistics
from pathlib import Path

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


def _clean_line(value):
    return " ".join(str(value or "").replace("\r", " ").split()).strip()


def _lines_to_text(lines):
    return "\n".join([_clean_line(item) for item in lines if _clean_line(item)]).strip()


def _title_from_filename(name):
    stem = Path(str(name or "")).stem
    cleaned = stem.replace("_", " ").replace("-", " ").strip()
    return cleaned or "Imported Module"


def _normalize_lesson(title, content, lesson_type, order):
    clean_title = _clean_line(title) or f"Lesson {order}"
    clean_content = str(content or "").strip()
    return {
        "title": clean_title,
        "content": clean_content,
        "type": lesson_type,
        "order": order,
    }


def _split_txt(raw_text):
    chunks = [chunk.strip() for chunk in str(raw_text or "").split("\n\n") if chunk.strip()]
    lessons = []
    for index, chunk in enumerate(chunks, start=1):
        first_line = _clean_line(chunk.split("\n")[0])
        title = first_line[:90] if first_line else f"Section {index}"
        lessons.append(_normalize_lesson(title, chunk, "paragraph", index))
    return lessons


def _extract_txt(raw_bytes):
    warnings = []
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            text = raw_bytes.decode(encoding)
            return _split_txt(text), warnings
        except Exception:
            continue
    text = raw_bytes.decode("utf-8", errors="ignore")
    warnings.append("TXT decoded with fallback encoding.")
    return _split_txt(text), warnings


def _extract_docx(raw_bytes):
    if docx is None:
        return [], ["python-docx is not installed. DOCX extraction unavailable."]

    document = docx.Document(io.BytesIO(raw_bytes))
    lessons = []
    warnings = []
    current_title = ""
    current_lines = []

    def flush(order_index):
        nonlocal current_title, current_lines
        if not current_title and not current_lines:
            return None
        text = _lines_to_text(current_lines)
        lesson = _normalize_lesson(current_title or f"Section {order_index}", text, "section", order_index)
        current_title = ""
        current_lines = []
        return lesson

    order = 1
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = str(getattr(paragraph.style, "name", "") or "").lower()
        is_heading = style_name.startswith("heading")
        if is_heading:
            previous = flush(order)
            if previous:
                lessons.append(previous)
                order += 1
            current_title = text
        else:
            current_lines.append(text)

    tail = flush(order)
    if tail:
        lessons.append(tail)

    if not lessons:
        fallback_text = _lines_to_text([p.text for p in document.paragraphs if p.text.strip()])
        if fallback_text:
            lessons = [_normalize_lesson("Imported DOCX", fallback_text, "paragraph", 1)]
        else:
            warnings.append("No readable paragraphs found in DOCX.")
    return lessons, warnings


def _extract_pptx(raw_bytes):
    if Presentation is None:
        return [], ["python-pptx is not installed. PPTX extraction unavailable."]

    deck = Presentation(io.BytesIO(raw_bytes))
    lessons = []
    warnings = []
    order = 1
    for slide in deck.slides:
        lines = []
        title = ""
        if getattr(slide.shapes, "title", None) and slide.shapes.title:
            title = _clean_line(slide.shapes.title.text)

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = _clean_line(shape.text_frame.text)
            if text:
                lines.append(text)

        if not title and lines:
            title = lines[0][:90]
        content = _lines_to_text(lines)
        if title or content:
            lessons.append(_normalize_lesson(title or f"Slide {order}", content, "slide", order))
            order += 1

    if not lessons:
        warnings.append("No readable slide text found in PPTX.")
    return lessons, warnings


def _extract_pdf(raw_bytes):
    if fitz is None:
        return [], ["PyMuPDF is not installed. PDF extraction unavailable."]

    lessons = []
    warnings = []
    order = 1

    with fitz.open(stream=raw_bytes, filetype="pdf") as document:
        all_font_sizes = []
        per_page_lines = []

        for page in document:
            raw = page.get_text("dict")
            lines = []
            for block in raw.get("blocks", []):
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    span_texts = []
                    span_sizes = []
                    for span in line.get("spans", []):
                        text = _clean_line(span.get("text", ""))
                        if text:
                            span_texts.append(text)
                            size = float(span.get("size", 0) or 0)
                            span_sizes.append(size)
                            all_font_sizes.append(size)
                    if span_texts:
                        lines.append(
                            {
                                "text": " ".join(span_texts).strip(),
                                "max_size": max(span_sizes) if span_sizes else 0,
                                "y": float(line.get("bbox", [0, 0, 0, 0])[1]),
                            }
                        )
            lines.sort(key=lambda item: item["y"])
            per_page_lines.append(lines)

        baseline = statistics.median(all_font_sizes) if all_font_sizes else 11
        heading_size = baseline + 1.8

        current_title = ""
        current_lines = []

        def flush_section():
            nonlocal current_title, current_lines, order
            text = _lines_to_text(current_lines)
            if not current_title and not text:
                return
            lessons.append(_normalize_lesson(current_title or f"Section {order}", text, "section", order))
            order += 1
            current_title = ""
            current_lines = []

        for page_lines in per_page_lines:
            for line in page_lines:
                text = line["text"]
                if not text:
                    continue
                likely_heading = (
                    line["max_size"] >= heading_size
                    and len(text) <= 120
                    and not text.endswith(".")
                )
                if likely_heading:
                    flush_section()
                    current_title = text
                else:
                    current_lines.append(text)
            current_lines.append("")

        flush_section()

    if not lessons:
        warnings.append("No readable text blocks found in PDF.")
    return lessons, warnings


def extract_module_structure(uploaded_file, module_title_override=""):
    """
    Parse file into a structured module payload.
    Output shape:
    {
      "module_title": "...",
      "lessons": [{"title": "...", "content": "...", "type": "...", "order": 1}, ...],
      "warnings": [...]
    }
    """
    file_name = str(getattr(uploaded_file, "name", "") or "")
    extension = Path(file_name).suffix.lower()
    raw_bytes = uploaded_file.read()
    if hasattr(uploaded_file, "seek"):
        uploaded_file.seek(0)

    warnings = []
    if extension == ".pptx":
        lessons, local_warnings = _extract_pptx(raw_bytes)
    elif extension == ".pdf":
        lessons, local_warnings = _extract_pdf(raw_bytes)
    elif extension == ".docx":
        lessons, local_warnings = _extract_docx(raw_bytes)
    elif extension == ".txt":
        lessons, local_warnings = _extract_txt(raw_bytes)
    elif extension in {".doc", ".ppt"}:
        lessons, local_warnings = [], [f"{extension} is supported for upload but extraction quality is limited."]
    else:
        lessons, local_warnings = [], [f"Unsupported extension: {extension}"]

    warnings.extend(local_warnings)

    normalized = []
    for index, lesson in enumerate(lessons, start=1):
        title = lesson.get("title") or f"Lesson {index}"
        content = str(lesson.get("content") or "").strip()
        if not content and not title:
            continue
        normalized.append(_normalize_lesson(title, content, lesson.get("type") or "paragraph", index))

    if not normalized:
        warnings.append("No structured lessons could be extracted from the file.")

    final_title = str(module_title_override or "").strip() or _title_from_filename(file_name)
    return {
        "module_title": final_title,
        "lessons": normalized,
        "warnings": warnings,
    }
